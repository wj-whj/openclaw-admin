#!/usr/bin/env python3
"""
generate_pptx.py - 将 JSON 内容转换为 .pptx 文件
用法: python3 generate_pptx.py <input.json> <output.pptx>
"""
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 风格配置
STYLES = {
    'tech': {
        'bg': RGBColor(0x0d, 0x0d, 0x1a),
        'title_color': RGBColor(0x18, 0xa0, 0xfb),
        'text_color': RGBColor(0xcc, 0xcc, 0xcc),
        'accent': RGBColor(0xa2, 0x59, 0xff),
        'title_font': 'Arial',
        'body_font': 'Arial',
    },
    'business': {
        'bg': RGBColor(0xff, 0xff, 0xff),
        'title_color': RGBColor(0x1a, 0x3a, 0x6b),
        'text_color': RGBColor(0x33, 0x33, 0x33),
        'accent': RGBColor(0x0d, 0x8d, 0xe3),
        'title_font': 'Calibri',
        'body_font': 'Calibri',
    },
    'minimal': {
        'bg': RGBColor(0xf8, 0xf8, 0xf8),
        'title_color': RGBColor(0x11, 0x11, 0x11),
        'text_color': RGBColor(0x44, 0x44, 0x44),
        'accent': RGBColor(0x55, 0x55, 0x55),
        'title_font': 'Helvetica Neue',
        'body_font': 'Helvetica Neue',
    },
}

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None,
                 font_name='Arial', align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    return txBox

def add_accent_bar(slide, style_cfg):
    """顶部装饰条"""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(10), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = style_cfg['accent']
    bar.line.fill.background()

def make_cover_slide(prs, slide_data, style_cfg):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, style_cfg['bg'])
    add_accent_bar(slide, style_cfg)

    # 主标题
    add_text_box(
        slide, slide_data.get('title', ''),
        left=0.8, top=2.2, width=8.4, height=1.4,
        font_size=40, bold=True,
        color=style_cfg['title_color'],
        font_name=style_cfg['title_font'],
        align=PP_ALIGN.CENTER
    )
    # 副标题
    subtitle = slide_data.get('subtitle', '')
    if subtitle:
        add_text_box(
            slide, subtitle,
            left=0.8, top=3.8, width=8.4, height=0.8,
            font_size=20, bold=False,
            color=style_cfg['text_color'],
            font_name=style_cfg['body_font'],
            align=PP_ALIGN.CENTER
        )

def make_content_slide(prs, slide_data, style_cfg):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, style_cfg['bg'])
    add_accent_bar(slide, style_cfg)

    # 页码
    page_num = slide_data.get('page', '')
    if page_num:
        add_text_box(
            slide, str(page_num),
            left=9.2, top=6.8, width=0.6, height=0.3,
            font_size=10, color=style_cfg['accent'],
            font_name=style_cfg['body_font']
        )

    # 标题
    add_text_box(
        slide, slide_data.get('title', ''),
        left=0.6, top=0.4, width=8.8, height=0.9,
        font_size=28, bold=True,
        color=style_cfg['title_color'],
        font_name=style_cfg['title_font']
    )

    # 分隔线（用细矩形模拟）
    line = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.35), Inches(8.8), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = style_cfg['accent']
    line.line.fill.background()

    # 要点列表
    points = slide_data.get('points', [])
    y = 1.55
    for point in points:
        if y > 6.2:
            break
        add_text_box(
            slide, f'• {point}',
            left=0.8, top=y, width=8.4, height=0.55,
            font_size=16, color=style_cfg['text_color'],
            font_name=style_cfg['body_font']
        )
        y += 0.6

    # 备注
    notes = slide_data.get('notes', '')
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

def generate(input_path, output_path):
    with open(input_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    style_name = data.get('style', 'tech')
    style_cfg = STYLES.get(style_name, STYLES['tech'])
    slides_data = data.get('slides', [])

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides_data):
        if i == 0:
            make_cover_slide(prs, slide_data, style_cfg)
        else:
            make_content_slide(prs, slide_data, style_cfg)

    prs.save(output_path)
    print(f'Generated: {output_path}')

if __name__ == '__main__':
    if len(sys.argv) != 3:
        print('Usage: python3 generate_pptx.py <input.json> <output.pptx>')
        sys.exit(1)
    generate(sys.argv[1], sys.argv[2])
